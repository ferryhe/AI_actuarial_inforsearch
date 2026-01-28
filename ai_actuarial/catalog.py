from __future__ import annotations

import os
import re
import json
import hashlib
import html as html_lib
import string
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from .storage import Storage
from .utils import load_category_config

# ---------------------------------------------------------------------------
# Catalog rules / keywords - Load from config file
# ---------------------------------------------------------------------------

# Load category configuration
try:
    _category_config = load_category_config()
    CATEGORY_RULES = _category_config.get("categories", {})
    AI_TERMS = _category_config.get("ai_filter_keywords", [])
    AI_KEYWORDS = _category_config.get("ai_keywords", [])
except FileNotFoundError as e:
    import logging
    logging.warning("Category config file not found, using default values: %s", e)
    # Fallback to hardcoded values if config load fails
    AI_TERMS = [
        "artificial intelligence",
        "machine learning",
        "deep learning",
        "neural",
        "llm",
        "large language model",
        "generative",
        "genai",
        "chatgpt",
        "openai",
        "transformer",
        "nlp",
    ]
    
    AI_KEYWORDS = [
        "artificial intelligence",
        "machine learning",
        "deep learning",
        "large language model",
        "generative ai",
        "llm",
        "genai",
        "chatgpt",
        "transformer",
        "neural network",
        "nlp",
    ]
    
    CATEGORY_RULES: dict[str, list[str]] = {
        "AI": [
            "artificial intelligence",
            "machine learning",
            "deep learning",
            "large language model",
            "llm",
            "generative ai",
            "neural network",
            "transformer model",
            "nlp",
            "chatgpt",
        ],
        "Regulation & Standards": [
            "ifrs",
            "solvency",
            "asop",
            "ias",
            "gaap",
            "naic",
            "standard",
            "compliance",
            "regulation",
        ],
        "Risk & Capital": [
            "erm",
            "risk",
            "capital",
            "stress",
            "scenario",
            "catastrophe",
            "reinsurance",
        ],
        "Pricing": [
            "pricing",
            "rate",
            "rating",
            "premium",
            "tariff",
        ],
        "Underwriting & Claims": [
            "underwriting",
            "uw",
            "risk selection",
            "appetite",
            "claim",
            "claims",
            "loss",
            "settlement",
        ],
        "Reserving": [
            "reserve",
            "reserving",
            "ibnr",
        ],
        "P&C": [
            "property",
            "casualty",
            "p&c",
            "auto",
            "general insurance",
        ],
        "Life": [
            "life",
            "annuity",
            "mortality",
            "longevity",
        ],
        "Health": [
            "health",
            "medical",
            "morbidity",
        ],
        "LTC / DI / CI": [
            "long term care",
            "ltc",
            "disability income",
            "di",
            "critical illness",
            "ci",
        ],
        "Data & Analytics": [
            "data",
            "analytics",
            "model",
            "modeling",
            "statistics",
            "forecast",
            "predictive",
            "regression",
            "time series",
            "governance",
        ],
        "Operations / Automation": [
            "automation",
            "workflow",
            "process",
            "rpa",
            "system",
            "implementation",
            "tooling",
        ],
        "Education / Events": [
            "webinar",
            "seminar",
            "conference",
            "agenda",
            "workshop",
            "training",
            "course",
            "syllabus",
            "lecture",
            "slides",
        ],
        "Investment / ALM": [
            "investment",
            "asset",
            "liability",
            "alm",
            "portfolio",
            "interest rate",
            "yield",
            "duration",
        ],
    }

_KEYBERT_MODEL = None
CATALOG_VERSION = "v2-keybert"

# ---------------------------------------------------------------------------
# Data model
# ---------------------------------------------------------------------------


@dataclass
class CatalogItem:
    source_site: str | None
    title: str | None
    original_filename: str | None
    url: str | None
    local_path: str | None
    keywords: list[str]
    summary: str
    category: str


# ---------------------------------------------------------------------------
# Text extraction (fast path + optional marker fallback) + lightweight caching
# ---------------------------------------------------------------------------

_SENTENCE_SPLIT = re.compile(r"(?<=[。！？.!?])\s+")


def detect_kind(path: Path) -> str:
    try:
        with open(path, "rb") as f:
            head = f.read(512)
    except OSError:
        return "unknown"

    if head.startswith(b"%PDF"):
        return "pdf"
    if head.startswith(b"\xff\xd8\xff"):
        return "jpeg"
    if head.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"

    stripped = head.lstrip()
    lowered = stripped[:256].lower()
    if lowered.startswith(b"<!doctype") or lowered.startswith(b"<html") or lowered.startswith(b"<!doc"):
        return "html"

    if head:
        null_ratio = head.count(b"\x00") / max(len(head), 1)
        if null_ratio <= 0.01:
            return "txt"

    return "unknown"


def _trim_semantic(text: str, max_chars: int) -> str:
    """Trim without chopping in the middle of a sentence when possible."""
    if not text:
        return ""
    if max_chars <= 0:
        return ""
    if len(text) <= max_chars:
        return text

    chunk = text[:max_chars]
    parts = _SENTENCE_SPLIT.split(chunk.strip())
    if len(parts) <= 1:
        return chunk.rstrip()

    # drop last potentially incomplete sentence
    trimmed = " ".join(parts[:-1]).rstrip()
    return trimmed if trimmed else chunk.rstrip()


def _looks_bad(text: str) -> bool:
    """Heuristic: detect failed/low-quality extraction (esp. scanned PDFs)."""
    if not text:
        return True
    if len(text) < 800:
        return True
    # Too many control chars often indicates extraction noise
    weird = sum(1 for c in text if ord(c) < 9 or (11 <= ord(c) < 32))
    if weird / max(len(text), 1) > 0.02:
        return True
    # Lots of repeated blank lines
    if text.count("\n\n\n") > 5:
        return True
    return False


def _cache_dir() -> Path:
    # Default under project root; override if you want
    base = os.getenv("CATALOG_CACHE_DIR", ".cache/catalog_extract")
    p = Path(base)
    p.mkdir(parents=True, exist_ok=True)
    return p


def _cache_key(path: Path, max_chars: int, extractor_version: str) -> str:
    # Keyed by absolute path + size + mtime + max_chars + extractor version
    try:
        st = path.stat()
        payload = f"{path.resolve()}|{st.st_size}|{int(st.st_mtime)}|{max_chars}|{extractor_version}"
    except Exception:
        payload = f"{path}|{max_chars}|{extractor_version}"
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _read_cache(path: Path, key: str) -> str | None:
    fp = _cache_dir() / f"{key}.json"
    if not fp.exists():
        return None
    try:
        obj = json.loads(fp.read_text(encoding="utf-8"))
        return obj.get("text") or None
    except Exception:
        return None


def _write_cache(key: str, text: str) -> None:
    fp = _cache_dir() / f"{key}.json"
    try:
        fp.write_text(
            json.dumps(
                {
                    "text": text,
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
    except Exception:
        # Cache is best-effort
        pass


def _read_pdf_fast(path: Path, max_chars: int, max_pages: int = 20) -> str:
    from pypdf import PdfReader
    import logging

    # Silence noisy parser warnings; errors are handled by caller.
    logging.getLogger("pypdf").setLevel(logging.ERROR)
    logging.getLogger("PyPDF2").setLevel(logging.ERROR)

    reader = PdfReader(str(path))
    out: list[str] = []
    used = 0

    for i, page in enumerate(reader.pages[:max_pages], start=1):
        t = page.extract_text() or ""
        t = re.sub(r"\n{3,}", "\n\n", t).strip()
        if not t:
            continue
        block = f"[p{i}]\n{t}\n"
        out.append(block)
        used += len(block)
        if used >= max_chars:
            break

    return _trim_semantic("\n".join(out), max_chars)


def _read_pdf_marker(path: Path, max_chars: int) -> str:
    """Slow path: use marker-pdf to get cleaner markdown, then flatten to text.

    Enable via env:
      PDF_USE_MARKER=1

    Note: marker API can differ by version; this function tries a couple patterns.
    """
    md: str | None = None

    # Try common marker entry points (best-effort)
    try:
        # marker>=? sometimes exposes PdfConverter like this
        from marker.converters.pdf import PdfConverter  # type: ignore

        converter = PdfConverter()
        md = converter.convert(str(path))  # type: ignore
    except Exception:
        md = None

    if md is None:
        try:
            # Some installs may provide a simple convert function
            from marker import convert  # type: ignore

            md = convert(str(path))  # type: ignore
        except Exception as e:
            raise RuntimeError("marker convert not available") from e

    # Flatten markdown -> plain text for keyword/category stage
    md = re.sub(r"`{3}.*?`{3}", " ", md, flags=re.S)  # code blocks
    md = re.sub(r"!\[.*?\]\(.*?\)", " ", md)  # images
    md = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", md)  # keep link text
    md = re.sub(r"^[#>\-*_]+\s*", "", md, flags=re.M)  # headings/quotes/bullets
    text = re.sub(r"\n{3,}", "\n\n", md).strip()
    return _trim_semantic(text, max_chars)


def _read_pdf(path: Path, max_chars: int) -> str:
    # Extraction version string: changes when you change logic
    extractor_version = "pdf_v2_fast_plus_marker_fallback"
    key = _cache_key(path, max_chars, extractor_version)
    cached = _read_cache(path, key)
    if cached is not None:
        return cached

    max_pages = int(os.getenv("PDF_MAX_PAGES", "20"))
    text = _read_pdf_fast(path, max_chars, max_pages=max_pages)

    if _looks_bad(text) and os.getenv("PDF_USE_MARKER") == "1":
        try:
            text = _read_pdf_marker(path, max_chars)
        except Exception:
            # keep fast result if marker fails
            pass

    _write_cache(key, text)
    return text


def _read_docx(path: Path, max_chars: int) -> str:
    import docx

    extractor_version = "docx_v2_para_plus_tables"
    key = _cache_key(path, max_chars, extractor_version)
    cached = _read_cache(path, key)
    if cached is not None:
        return cached

    doc = docx.Document(str(path))
    parts: list[str] = []

    # Paragraphs
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)

    # Tables (often critical for business/regulatory docs)
    if os.getenv("DOCX_TABLES_DISABLE") != "1":
        for table in doc.tables:
            for row in table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(" ".join(c.split()) for c in cells))

    text = _trim_semantic("\n".join(parts), max_chars)
    _write_cache(key, text)
    return text


def _read_pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation

    extractor_version = "pptx_v2_slide_boundaries"
    key = _cache_key(path, max_chars, extractor_version)
    cached = _read_cache(path, key)
    if cached is not None:
        return cached

    prs = Presentation(str(path))
    parts: list[str] = []
    used = 0

    for s_idx, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_texts.append(t)

        if not slide_texts:
            continue

        block = f"[slide {s_idx}]\n" + "\n".join(slide_texts) + "\n"
        parts.append(block)
        used += len(block)
        if used >= max_chars:
            break

    text = _trim_semantic("\n".join(parts), max_chars)
    _write_cache(key, text)
    return text


def extract_text(path: Path, max_chars: int = 20000) -> str:
    """Extract text for lightweight keyword/category + heuristic summary.

    Notes:
    - This is optimized for speed and robustness, not perfect layout fidelity.
    - PDF uses fast pypdf extraction, with optional marker fallback when enabled.
    - Extraction results are cached to avoid repeatedly re-parsing the same file.
    """
    if not path.exists():
        return ""
    kind = detect_kind(path)
    if kind == "pdf":
        try:
            return _read_pdf(path, max_chars)
        except Exception as e:
            raise ValueError(
                f"corrupt pdf ({kind}): {path} ({type(e).__name__}: {e})"
            ) from e
    if kind == "html":
        return _read_html(path, max_chars)
    if kind in {"jpeg", "png"}:
        return _read_image_ocr(path, max_chars, kind)
    if kind == "txt":
        return _read_txt(path, max_chars)

    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _read_docx(path, max_chars)
    if suffix == ".pptx":
        return _read_pptx(path, max_chars)
    raise ValueError(f"unsupported file type ({kind}): {path}")


def _read_html(path: Path, max_chars: int) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator=" ")
        text = re.sub(r"\s+", " ", text).strip()
        return _trim_semantic(text, max_chars)
    except Exception:
        # Fallback: strip tags + unescape entities
        text = re.sub(r"(?s)<(script|style).*?>.*?</\\1>", " ", raw, flags=re.I)
        text = re.sub(r"<[^>]+>", " ", text)
        text = html_lib.unescape(text)
        text = re.sub(r"\s+", " ", text).strip()
        return _trim_semantic(text, max_chars)


def _read_txt(path: Path, max_chars: int) -> str:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return _trim_semantic(text, max_chars)


def _read_image_ocr(path: Path, max_chars: int, kind: str) -> str:
    try:
        import pytesseract  # type: ignore
    except Exception as e:
        raise ValueError(f"ocr unavailable for {kind}: install pytesseract") from e
    try:
        from PIL import Image  # type: ignore
    except Exception as e:
        raise ValueError(f"ocr unavailable for {kind}: install pillow") from e

    img = Image.open(path)
    try:
        text = pytesseract.image_to_string(img)
    except Exception as e:
        raise ValueError(f"ocr failed for {kind}: {path} ({type(e).__name__}: {e})") from e
    text = re.sub(r"\s+", " ", text).strip()
    return _trim_semantic(text, max_chars)


# ---------------------------------------------------------------------------
# Summarization / keyword extraction / categorization (existing logic)
# ---------------------------------------------------------------------------


def _split_sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+", text.strip())
    return [p.strip() for p in parts if len(p.strip()) > 20]


def summarize(text: str, keywords: list[str], max_sentences: int = 4) -> str:
    sentences = _split_sentences(text)
    if not sentences:
        return ""
    if not keywords:
        return " ".join(sentences[:max_sentences])
    keyset = [k.lower() for k in keywords]
    scored = []
    for s in sentences:
        sl = s.lower()
        score = sum(1 for k in keyset if k in sl)
        scored.append((score, s))
    scored.sort(key=lambda x: x[0], reverse=True)
    picked = [s for _, s in scored[:max_sentences]]
    return " ".join(picked)


def _ai_hit(title: str | None, keywords: list[str]) -> bool:
    """Check if document has explicit AI/ML keywords - using word-boundary matching."""
    hay = (title or "") + " " + " ".join(keywords)
    hay = hay.lower()
    words = set(re.findall(r'\b\w+\b', hay))
    
    # Only match full words/phrases to avoid false positives
    ai_terms = ["artificial intelligence", "machine learning", "deep learning", 
                "large language model", "llm", "generative ai", "neural network", "nlp"]
    
    for term in ai_terms:
        if " " in term:
            # Multi-word: check phrase in hay
            if term in hay:
                return True
        else:
            # Single word: check in word set
            if term in words:
                return True
    
    return False


def is_ai_related(text: str, keywords: list[str], title: str | None = None) -> bool:
    # Keep broader AI detection for optional ai_only filtering
    hay = (title or "") + " " + text
    hay = hay.lower()
    for term in AI_TERMS:
        if term in hay:
            return True
    for k in keywords:
        if k.lower() in hay:
            return True
    return False


def extract_keywords(text: str, title: str | None = None, top_n: int = 8) -> list[str]:
    if not text and not title:
        return []
    if os.getenv("KEYBERT_DISABLE") == "1":
        return _light_keywords(text, title, top_n)
    global _KEYBERT_MODEL
    try:
        from keybert import KeyBERT

        if _KEYBERT_MODEL is None:
            _KEYBERT_MODEL = KeyBERT()
        base = (title or "") + "\n" + text
        kw = _KEYBERT_MODEL.extract_keywords(base, top_n=top_n, stop_words="english")
        return [k for k, _ in kw]
    except Exception:
        return _light_keywords(text, title, top_n)


def _light_keywords(text: str, title: str | None, top_n: int) -> list[str]:
    combined = ((title or "") + "\n" + (text or "")).lower()
    combined = combined.replace("\u2013", "-").replace("\u2014", "-")

    specials = []
    for pat in [r"\bifrs\s*17\b", r"\bcovid[-\s]*19\b", r"\bsolvency\s*ii\b"]:
        m = re.findall(pat, combined, flags=re.IGNORECASE)
        specials.extend(m)

    def tokens(s: str) -> list[str]:
        return re.findall(r"[a-zA-Z][a-zA-Z0-9\-]*", s.lower())

    text_tokens = tokens(text or "")
    title_tokens = tokens(title or "")

    freq: dict[str, int] = {}
    for t in text_tokens:
        if len(t) < 3:
            continue
        if t.isdigit():
            continue
        freq[t] = freq.get(t, 0) + 1

    bigrams: dict[str, int] = {}
    for a, b in zip(text_tokens, text_tokens[1:]):
        if a.isdigit() and b.isdigit():
            continue
        phrase = f"{a} {b}"
        bigrams[phrase] = bigrams.get(phrase, 0) + 1

    title_set = set(title_tokens)
    title_bigrams = set(f"{a} {b}" for a, b in zip(title_tokens, title_tokens[1:]))

    scored: dict[str, int] = {}
    for t, c in freq.items():
        scored[t] = scored.get(t, 0) + c + (3 if t in title_set else 0)
    for t, c in bigrams.items():
        scored[t] = scored.get(t, 0) + c + (4 if t in title_bigrams else 0)

    for s in specials:
        s_norm = re.sub(r"\s+", " ", s.strip().lower())
        scored[s_norm] = scored.get(s_norm, 0) + 10

    ranked = sorted(scored.items(), key=lambda x: x[1], reverse=True)
    out: list[str] = []
    for term, _ in ranked:
        if term.isdigit():
            continue
        if term not in out:
            out.append(term)
        if len(out) >= top_n:
            break
    return out


def categorize(title: str | None, text: str, keywords: list[str]) -> str:
    """Categorize document using word-boundary matching for accuracy."""
    hay = (title or "") + " " + text + " " + " ".join(keywords)
    hay = hay.lower()
    
    # Tokenize to word list for word-boundary matching
    words = re.findall(r'\b\w+\b', hay)
    word_set = set(words)
    
    matches: list[tuple[str, int]] = []
    for cat, terms in CATEGORY_RULES.items():
        score = 0
        for t in terms:
            # Use word-boundary matching: check if term matches complete words
            term_words = t.split()
            if len(term_words) == 1:
                # Single word: exact match in word set
                if t in word_set:
                    score += 1
            else:
                # Multi-word term: check if sequence exists
                pattern = r'\b' + r'\s+'.join(re.escape(w) for w in term_words) + r'\b'
                if re.search(pattern, hay):
                    score += 1
        
        if score > 0:
            matches.append((cat, score))
    
    if not matches:
        return "Other"
    
    matches.sort(key=lambda x: x[1], reverse=True)
    cats = [c for c, _ in matches]
    
    # Only force AI to front if it was already in matches (don't artificially boost)
    # if _ai_hit(title, keywords) and "AI" not in cats:
    #     cats = ["AI"] + cats
    
    return "; ".join(cats[:3])


def _fallback_keywords(text: str, top_n: int) -> list[str]:
    try:
        from sklearn.feature_extraction.text import TfidfVectorizer

        vectorizer = TfidfVectorizer(
            stop_words="english",
            ngram_range=(1, 2),
            max_features=50,
        )
        tfidf = vectorizer.fit_transform([text])
        scores = tfidf.toarray()[0]
        terms = vectorizer.get_feature_names_out()
        ranked = sorted(zip(terms, scores), key=lambda x: x[1], reverse=True)
        return [t for t, _ in ranked[:top_n]]
    except Exception:
        return []


# ---------------------------------------------------------------------------
# Catalog build / output
# ---------------------------------------------------------------------------


def build_catalog(
    storage: Storage,
    site_filter: str | None,
    limit: int,
    ai_only: bool = False,
    offset: int = 0,
) -> list[CatalogItem]:
    rows = storage.iter_files(
        site_filter=site_filter,
        limit=limit,
        offset=offset,
        require_local_path=True,
    )
    items: list[CatalogItem] = []
    filters: list[str] = []
    if site_filter:
        filters = [s.strip().lower() for s in site_filter.split(",") if s.strip()]
    for row in rows:
        site = row.get("source_site")
        if filters:
            if not site or not any(f in site.lower() for f in filters):
                continue
        path = row.get("local_path")
        if not path:
            continue
        try:
            text = extract_text(Path(path))
        except Exception:
            continue
        if not text:
            continue
        title = row.get("title")
        keywords = extract_keywords(text, title=title)
        if ai_only and not is_ai_related(text, keywords, title=title):
            continue
        summary = summarize(text, keywords)
        category = categorize(title, text, keywords)
        items.append(
            CatalogItem(
                source_site=site,
                title=title,
                original_filename=row.get("original_filename"),
                url=row.get("url"),
                local_path=path,
                keywords=keywords,
                summary=summary,
                category=category,
            )
        )
        if limit is not None and len(items) >= limit:
            break
    return items


def build_catalog_batch(
    storage: Storage,
    site_filter: str | None,
    row_limit: int,
    ai_only: bool = False,
    offset: int = 0,
) -> tuple[list[CatalogItem], int]:
    rows = storage.iter_files(
        site_filter=site_filter,
        limit=row_limit,
        offset=offset,
        require_local_path=True,
    )
    items: list[CatalogItem] = []
    filters: list[str] = []
    if site_filter:
        filters = [s.strip().lower() for s in site_filter.split(",") if s.strip()]
    for row in rows:
        site = row.get("source_site")
        if filters:
            if not site or not any(f in site.lower() for f in filters):
                continue
        path = row.get("local_path")
        if not path:
            continue
        try:
            text = extract_text(Path(path))
        except Exception:
            continue
        if not text:
            continue
        title = row.get("title")
        keywords = extract_keywords(text, title=title)
        if ai_only and not is_ai_related(text, keywords, title=title):
            continue
        summary = summarize(text, keywords)
        category = categorize(title, text, keywords)
        items.append(
            CatalogItem(
                source_site=site,
                title=title,
                original_filename=row.get("original_filename"),
                url=row.get("url"),
                local_path=path,
                keywords=keywords,
                summary=summary,
                category=category,
            )
        )
    return items, len(rows)


def write_catalog_jsonl(path: Path, items: Iterable[dict]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "a", encoding="utf-8") as f:
        for item in items:
            f.write(json.dumps(item, ensure_ascii=False) + "\n")


def build_catalog_incremental(
    storage: Storage,
    site_filter: str | None,
    limit: int | None,
    offset: int = 0,
    ai_only: bool = False,
    pipeline_version: str | None = None,
    retry_errors: bool = False,
    batch_size: int = 50,
) -> list[dict]:
    effective_version = pipeline_version or CATALOG_VERSION
    rows = storage.iter_files(
        site_filter=site_filter,
        limit=limit,
        offset=offset,
        require_local_path=True,
        only_changed=True,
        extractor_version=effective_version,
        include_errors=retry_errors,
    )
    results: list[dict] = []
    filters: list[str] = []
    processed = 0
    skipped = 0
    errors = 0
    if site_filter:
        filters = [s.strip().lower() for s in site_filter.split(",") if s.strip()]
    if batch_size <= 0:
        batch_size = len(rows) if rows else 1
    for start in range(0, len(rows), batch_size):
        batch = rows[start:start + batch_size]
        with storage.transaction():
            for row in batch:
                site = row.get("source_site")
                if filters:
                    if not site or not any(f in site.lower() for f in filters):
                        skipped += 1
                        continue
                path = row.get("local_path")
                if not path:
                    skipped += 1
                    continue
                try:
                    text = extract_text(Path(path))
                    if not text:
                        # Keep existing skip behavior for empty extraction results.
                        skipped += 1
                        continue
                    title = row.get("title")
                    keywords = extract_keywords(text, title=title)
                    if ai_only and not is_ai_related(text, keywords, title=title):
                        skipped += 1
                        continue
                    summary = summarize(text, keywords)
                    category = categorize(title, text, keywords)
                    item = {
                        "source_site": site,
                        "title": title,
                        "original_filename": row.get("original_filename"),
                        "url": row.get("url"),
                        "local_path": path,
                        "keywords": keywords,
                        "summary": summary,
                        "category": category,
                        "sha256": row.get("sha256"),
                    }
                    ts = storage.now()
                    storage.upsert_catalog_item(
                        item,
                        pipeline_version=effective_version,
                        status="ok",
                        error=None,
                        processed_at=ts,
                    )
                    results.append(item)
                    processed += 1
                except Exception as e:
                    ts = storage.now()
                    storage.upsert_catalog_item(
                        {"url": row.get("url"), "sha256": row.get("sha256")},
                        pipeline_version=effective_version,
                        status="error",
                        error=str(e),
                        processed_at=ts,
                    )
                    errors += 1
                    continue
    if processed or skipped or errors:
        print(
            f"catalog_incremental: processed={processed} skipped={skipped} errors={errors}"
        )
    return results


def write_catalog_md(path: Path, items: Iterable[CatalogItem], append: bool = False) -> None:
    headers = [
        "category",
        "source_site",
        "title",
        "original_filename",
        "keywords",
        "summary",
        "url",
        "local_path",
    ]
    path.parent.mkdir(parents=True, exist_ok=True)
    mode = "a" if append else "w"
    with open(path, mode, encoding="utf-8") as f:
        if not append:
            f.write("| " + " | ".join(headers) + " |\n")
            f.write("| " + " | ".join(["---"] * len(headers)) + " |\n")
        for item in items:
            values = [
                item.category,
                item.source_site or "",
                item.title or "",
                item.original_filename or "",
                ", ".join(item.keywords),
                item.summary,
                item.url or "",
                item.local_path or "",
            ]
            safe = [v.replace("|", " ") for v in values]
            f.write("| " + " | ".join(safe) + " |\n")
